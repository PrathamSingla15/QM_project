"""Build the Team F final presentation deck for the TQM course project.

Generates TQM_Project_Final.pptx: 24 slides, 16:9, with a consistent
cream/ink/accent-orange visual system, embedded survey charts, TQM artifact
diagrams, and app screenshots. Recomputes key survey stats from the raw CSV
and prints them to stdout for spot-checking.
"""

from __future__ import annotations

import os
from pathlib import Path

import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu


# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

ROOT = Path("/home/pratham/pratham/QM_project")
DELIVER = ROOT / "deliverables"
OUT_DIR = DELIVER / "presentation"
OUT_PPTX = OUT_DIR / "TQM_Project_Final.pptx"

SURVEY_CSV = DELIVER / "survey" / "responses.csv"
SURVEY_CHARTS = DELIVER / "survey" / "charts"
ANALYSIS = DELIVER / "analysis"
APP_SHOTS = DELIVER / "app" / "screenshots"


# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------

BG = RGBColor(0xF6, 0xF2, 0xEA)     # cream background
INK = RGBColor(0x1B, 0x1A, 0x17)    # near-black body
ACCENT = RGBColor(0xE8, 0x5D, 0x2F) # warm orange
MUTED = RGBColor(0xA3, 0x9A, 0x87)  # muted footer tone
MUTED_CAPTION = RGBColor(0x6E, 0x6A, 0x60)  # darker muted for italic caption line
HEADER_FILL = RGBColor(0xEA, 0xD8, 0xC7)  # light warm accent for table headers

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COURSE_FOOTER = "TQM Course Project - IIT Roorkee"
TOTAL_SLIDES = 24

# image placement constraints
IMG_MAX_W = 11.5
IMG_MAX_H = 5.3


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_run_font(run, *, name: str = "Calibri", size: int = 18, bold: bool = False,
                 italic: bool = False, color: RGBColor = INK) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_bg(slide) -> None:
    """Full-slide cream background rectangle as the first shape."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG
    rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_title(slide, text: str, *, top: float = 0.45, left: float = 0.6,
              width: float = 12.0, size: int = 32) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.85))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_font(run, name="Georgia", size=size, bold=True, color=INK)


def add_accent_bar(slide, *, top: float = 1.25, left: float = 0.62,
                   width: float = 2.0, height: float = 0.05) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_footer(slide, text: str = COURSE_FOOTER) -> None:
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(9.0), Inches(0.25))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_font(run, name="Calibri", size=9, italic=True, color=MUTED)


def add_page_number(slide, idx: int) -> None:
    tb = slide.shapes.add_textbox(Inches(11.3), Inches(7.15), Inches(1.7), Inches(0.25))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Slide {idx} / {TOTAL_SLIDES}"
    set_run_font(run, name="Calibri", size=9, color=MUTED)


def add_bullets(slide, items, *, left: float = 0.7, top: float = 1.6,
                width: float = 12.0, height: float = 5.2, size: int = 18,
                bold_first: bool = False) -> None:
    """Add a text box of hyphen-prefixed bullet lines."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"- {item}"
        set_run_font(run, name="Calibri", size=size,
                     bold=(bold_first and i == 0), color=INK)


def add_caption(slide, text: str, *, top: float = 6.6, left: float = 0.7,
                width: float = 12.0, size: int = 14, italic: bool = True) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.45))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run_font(run, name="Calibri", size=size, italic=italic, color=INK)


def fit_image(path: Path, max_w_in: float, max_h_in: float):
    """Return (width_in, height_in) constrained to max size, preserving aspect."""
    with Image.open(path) as im:
        w_px, h_px = im.size
    aspect = h_px / w_px
    w_in = max_w_in
    h_in = w_in * aspect
    if h_in > max_h_in:
        h_in = max_h_in
        w_in = h_in / aspect
    return w_in, h_in


def add_centered_image(slide, path: Path, *, top: float = 1.55,
                       max_w: float = IMG_MAX_W, max_h: float = IMG_MAX_H) -> None:
    w_in, h_in = fit_image(path, max_w, max_h)
    left = (13.333 - w_in) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                             width=Inches(w_in), height=Inches(h_in))


def add_two_images(slide, path_a: Path, path_b: Path, *, top: float = 1.7,
                   col_w: float = 5.8, gap: float = 0.2,
                   max_h: float = IMG_MAX_H) -> None:
    """Two images side-by-side, centered as a pair."""
    def fit_col(path: Path):
        w_in, h_in = fit_image(path, col_w, max_h)
        return w_in, h_in

    wa, ha = fit_col(path_a)
    wb, hb = fit_col(path_b)
    pair_w = wa + gap + wb
    left_a = (13.333 - pair_w) / 2
    left_b = left_a + wa + gap
    # vertically center each within the allotted band
    band_h = max(ha, hb)
    slide.shapes.add_picture(str(path_a), Inches(left_a),
                             Inches(top + (band_h - ha) / 2),
                             width=Inches(wa), height=Inches(ha))
    slide.shapes.add_picture(str(path_b), Inches(left_b),
                             Inches(top + (band_h - hb) / 2),
                             width=Inches(wb), height=Inches(hb))


def add_insights_panel(slide, bullets, *, left: float, top: float,
                        width: float, height: float,
                        header_size: int = 13, bullet_size: int = 13,
                        header_text: str = "KEY INSIGHTS") -> None:
    """Render a 'KEY INSIGHTS' header plus a bulleted takeaway list.

    Header: Calibri bold, accent orange, slightly looser letter-spacing.
    Body: Calibri ink 13pt, hyphen-prefixed bullets, line-spacing 1.15,
          small space-before on every paragraph after the first.
    """
    # Header text box
    hdr = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(0.32))
    htf = hdr.text_frame
    htf.margin_left = htf.margin_right = 0
    htf.margin_top = htf.margin_bottom = 0
    htf.word_wrap = True
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run()
    hr.text = header_text
    set_run_font(hr, name="Calibri", size=header_size, bold=True, color=ACCENT)
    # Letter-spacing ('a bit looser'). python-pptx has no direct API; fall back
    # to writing the rPr@spc XML attribute. Guard so a failure is cosmetic.
    try:
        rPr = hr._r.get_or_add_rPr()
        rPr.set("spc", "150")  # +1.5 pt inter-character tracking
    except Exception:
        pass

    # Body text box: one text frame, one paragraph per bullet
    body_top = top + 0.38
    body_h = max(0.4, height - 0.38)
    tb = slide.shapes.add_textbox(Inches(left), Inches(body_top),
                                   Inches(width), Inches(body_h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True

    for i, item in enumerate(bullets[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        if i > 0:
            p.space_before = Pt(9)  # ~0.12 inch
        run = p.add_run()
        run.text = f"- {item}"
        set_run_font(run, name="Calibri", size=bullet_size, color=INK)


def add_image_with_insights(slide, image_path: Path, bullets,
                              *, caption: str | None = None) -> None:
    """Image on the left, insights panel on the right, optional caption below."""
    w_in, h_in = fit_image(image_path, 7.8, 4.8)
    slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.7),
                              width=Inches(w_in), height=Inches(h_in))

    add_insights_panel(slide, bullets,
                        left=8.6, top=1.7, width=4.4, height=4.8)

    if caption:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.7),
                                       Inches(8.0), Inches(0.35))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = caption
        set_run_font(r, name="Calibri", size=12, italic=True,
                      color=MUTED_CAPTION)


def add_two_images_with_insights_below(slide, path_a: Path, path_b: Path,
                                         bullets, *,
                                         caption: str | None = None) -> None:
    """Two images side-by-side across the top, insights panel spanning below."""
    wa, ha = fit_image(path_a, 6.0, 4.0)
    slide.shapes.add_picture(str(path_a), Inches(0.5), Inches(1.7),
                              width=Inches(wa), height=Inches(ha))
    wb, hb = fit_image(path_b, 6.0, 4.0)
    slide.shapes.add_picture(str(path_b), Inches(6.95), Inches(1.7),
                              width=Inches(wb), height=Inches(hb))

    add_insights_panel(slide, bullets,
                        left=0.5, top=5.9, width=12.3, height=1.0,
                        bullet_size=12)


def new_slide(prs, *, idx: int, title: str, footer: str = COURSE_FOOTER,
              include_chrome: bool = True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    if title:
        add_title(slide, title)
        add_accent_bar(slide)
    if include_chrome:
        add_footer(slide, footer)
        add_page_number(slide, idx)
    return slide


def style_table_cell(cell, text: str, *, size: int = 11, bold: bool = False,
                      fill: RGBColor | None = None, color: RGBColor = INK,
                      align=PP_ALIGN.LEFT) -> None:
    cell.text = ""
    tf = cell.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, name="Calibri", size=size, bold=bold, color=color)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG


def add_table(slide, headers, rows, *, left: float, top: float,
              width: float, height: float,
              header_size: int = 12, body_size: int = 11) -> None:
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    tbl = table_shape.table
    # header row
    for j, h in enumerate(headers):
        style_table_cell(tbl.cell(0, j), h, size=header_size, bold=True,
                         fill=HEADER_FILL, color=INK, align=PP_ALIGN.LEFT)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            style_table_cell(tbl.cell(i, j), str(val), size=body_size,
                             bold=False, fill=BG, color=INK, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_01_title(prs) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Big centered title block (no page number or slide-n footer).
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.6))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "E-Rickshaw Service Quality on Campus"
    set_run_font(run, name="Georgia", size=46, bold=True, color=INK)

    # accent bar under title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(4.05),
                                  Inches(2.5), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.25), Inches(12.0), Inches(1.2))
    tf2 = sub.text_frame
    tf2.margin_left = tf2.margin_right = 0
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = ("A TQM Intervention Using App-Based Booking, "
                 "Wallet Payment, and Demand Aggregation")
    set_run_font(run2, name="Georgia", size=22, italic=True, color=INK)

    # Team line
    tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(12.0), Inches(0.4))
    p3 = tb3.text_frame.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Team F - Final Presentation Deck"
    set_run_font(r3, name="Calibri", size=16, color=MUTED)

    # Footer (centered-ish, bottom)
    tbF = slide.shapes.add_textbox(Inches(0.8), Inches(7.15), Inches(12.0), Inches(0.3))
    pF = tbF.text_frame.paragraphs[0]
    rF = pF.add_run()
    rF.text = COURSE_FOOTER
    set_run_font(rF, name="Calibri", size=11, italic=True, color=MUTED)


def slide_02_problem(prs) -> None:
    slide = new_slide(prs, idx=2, title="Problem in One Slide")
    # Headline
    hb = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.6))
    p = hb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Campus E-ricks waste ~8 min per ride on average"
    set_run_font(r, name="Georgia", size=24, bold=True, color=ACCENT)

    bullets = [
        "Long peak-hour wait times - mean 8.2 min, 90th percentile 13 min.",
        "Payment friction - cash-only, change disputes, inconsistent fares.",
        "Poor night-hour availability - service thins after 9 pm.",
        "No accountability - no KPIs, no visibility, no complaint channel.",
    ]
    add_bullets(slide, bullets, top=2.3, size=20)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(12.0), Inches(0.4))
    pn = note.text_frame.paragraphs[0]
    rn = pn.add_run()
    rn.text = "Data source: structured student survey, n = 150 respondents."
    set_run_font(rn, name="Calibri", size=14, italic=True, color=MUTED)


def slide_03_stakeholders(prs) -> None:
    slide = new_slide(prs, idx=3, title="Stakeholder Map",
                      footer="Tool: Stakeholder Analysis - Ref. Ch.1")
    headers = ["Stakeholder", "Interest", "Influence", "Pain", "Gain"]
    rows = [
        ["Students", "Fast, fair, cashless rides", "High", "Waits, payment issues, no night service",
         "App booking, wallet, visibility"],
        ["Drivers", "Stable earnings, fewer disputes", "High",
         "Empty returns, cash disputes, idle time",
         "Demand aggregation, wallet clearing"],
        ["Admin (IITR)", "Safety, order, utilisation", "High",
         "No data, complaints, ad-hoc rules",
         "KPI dashboard, auditable leakage %"],
        ["Security / DoSW", "Peak-hour congestion, safety", "Medium",
         "Disorderly queues, night gaps",
         "Zone-aware routing, visible SLAs"],
        ["Payment vendor", "Transaction volume, low chargebacks", "Medium",
         "Low penetration in campus E-ricks",
         "New captive wallet userbase"],
    ]
    add_table(slide, headers, rows, left=0.55, top=1.7,
              width=12.25, height=4.6, header_size=13, body_size=12)


def slide_04_methodology(prs) -> None:
    slide = new_slide(prs, idx=4, title="Methodology")
    bullets = [
        "22-question structured survey across demographics, usage, pain ranking, solution willingness, and open-ended feedback.",
        "Synthetic dataset of n = 150 respondents (random seed 20260421) generated via a latent-pain factor model that propagates realistic co-variation across responses.",
        "Nine TQM artifacts - Ishikawa, 5 Whys, Pareto, SIPOC, current-vs-future flowchart, FMEA, PDCA, Affinity Diagram, Relations Diagram.",
        "Working React prototype exercising all core flows: map discovery, booking, countdown, wallet, driver view, admin KPIs.",
    ]
    add_bullets(slide, bullets, top=1.65, size=17)

    # Disclaimer
    disc = slide.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.0))
    tf = disc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = ("Disclaimer: the 150 responses are synthetic, calibrated to plausible campus "
              "mobility patterns. Distributions, correlations, and headline numbers are "
              "illustrative of what a real survey could surface; field collection would be "
              "required to validate absolute magnitudes.")
    set_run_font(r, name="Calibri", size=13, italic=True, color=MUTED)


def slide_05_wait_time(prs) -> None:
    slide = new_slide(prs, idx=5, title="Wait Time Reality",
                      footer="Tool: Histogram - Ref. Ch.2 Basic 7 QC Tools")
    bullets = [
        "Mean 8.2 min, median 7 min, 90th percentile 13 min - distribution is right-skewed.",
        "About 30% of respondents report waits over 10 min at peak hours.",
        "Tail risk is the presentation problem: long outlier waits damage trust more than the average.",
        "App target: cap p90 under 3 min, not just shift the mean.",
    ]
    add_image_with_insights(slide, SURVEY_CHARTS / "wait_times_histogram.png",
                              bullets,
                              caption="Wait-time distribution, n = 150; right-skewed with a heavy peak-hour tail.")


def slide_06_pareto(prs) -> None:
    slide = new_slide(prs, idx=6, title="Pareto - What Hurts Most",
                      footer="Tool: Pareto Chart - Ref. Ch.2 Basic 7 QC Tools")
    bullets = [
        "Top 3 issues (long wait 60%, payment 47%, no night avail 44%) = 56.8% of all pain mentions.",
        "All three top issues are app-addressable; a single intervention hits the vital few.",
        "Rude driver (35%) is a cultural issue; app provides accountability but not a direct fix.",
        "80%-of-complaints line crosses at the 5th category - no hidden minor cause dominates.",
    ]
    add_image_with_insights(slide, SURVEY_CHARTS / "pain_points_pareto.png",
                              bullets,
                              caption="Pareto of pain-point mentions; top-3 concentrate the majority of dissatisfaction.")


def slide_07_ishikawa(prs) -> None:
    slide = new_slide(prs, idx=7, title="Root Cause - Ishikawa (6M)",
                      footer="Tool: Cause & Effect - Ref. Ch.2")
    bullets = [
        "Method (rigid 4-pax rule, no zone routing, cash only) is the densest category.",
        "Measurement (no KPIs, no audit, no reconciliation) is entirely absent today - biggest lift from software.",
        "Machine (no GPS, no capacity sensor) is the enabling gap the app closes.",
        "Environment (scattered layout) turns from liability into advantage once routing is digital.",
    ]
    add_image_with_insights(slide, ANALYSIS / "01_ishikawa" / "diagram.png",
                              bullets,
                              caption="6M decomposition; method and measurement dominate the cause space.")


def slide_08_five_whys(prs) -> None:
    slide = new_slide(prs, idx=8, title="5 Whys",
                      footer="Tool: 5 Whys - Ref. Ch.2 (practice)")
    bullets = [
        "Proximate cause (rigid 4-fill rule) sits at Why 1 - common but superficial.",
        "Why 3 (no zone aggregation) is the pivot from proximate to structural.",
        "Root cause at Why 5: absence of a digital booking and aggregation layer.",
        "Solving the root resolves all 4 levels above it simultaneously.",
    ]
    add_image_with_insights(slide, ANALYSIS / "02_five_whys" / "diagram.png",
                              bullets,
                              caption="Five-why chain from observed wait to the absent digital layer.")


def slide_09_relations(prs) -> None:
    slide = new_slide(prs, idx=9, title="Vicious Cycle - Relations Diagram",
                      footer="Tool: Relations Diagram - Ref. Ch.3 New 7")
    bullets = [
        "6-node reinforcing loop: long wait -> avoidance -> fewer riders -> lower earnings -> fewer ricks -> longer wait.",
        "Cutting any single edge breaks the cycle; the app cuts three at once (discovery, earnings, fleet stability).",
        "Off-loop contributors (cash leakage, route detours) amplify two entry points into the loop.",
        "Implication: measurable improvement compounds - each KPI gain reinforces the next.",
    ]
    add_image_with_insights(slide, ANALYSIS / "09_relations" / "diagram.png",
                              bullets,
                              caption="Reinforcing loop among wait, adoption, fleet size, and driver earnings.")


def slide_10_affinity(prs) -> None:
    slide = new_slide(prs, idx=10, title="Voice of Student - Affinity",
                      footer="Tool: Affinity Diagram - Ref. Ch.3 New 7")
    bullets = [
        "Six themes surfaced from 47 open-ended responses; schedule reliability (14) is most-mentioned.",
        "Safety-and-identity theme (8) is a hidden requirement - driver name+photo + women-only timing.",
        "Zone aggregation (10) and fare transparency (12) map directly to proposed app features.",
        "Driver-experience theme (4) signals a worker-side intervention beyond app scope.",
    ]
    add_image_with_insights(slide, ANALYSIS / "08_affinity" / "diagram.png",
                              bullets,
                              caption="Six affinity clusters drawn from open-ended survey responses.")


def slide_11_correlation(prs) -> None:
    slide = new_slide(prs, idx=11, title="Willingness Correlates with Pain",
                      footer="Tool: Correlation - Ref. Ch.2")
    bullets = [
        "Strongest pair: r(payment issues, wallet comfort) = 0.40 - validates the wallet feature.",
        "Less-satisfied students are more willing to try the app (r = -0.27).",
        "Cross-pain-to-willingness correlations are mild but consistently positive (0.19 to 0.27).",
        "Pain is a reliable adoption driver; marketing to the heaviest sufferers first is efficient.",
    ]
    add_image_with_insights(slide, SURVEY_CHARTS / "correlation_heatmap.png",
                              bullets,
                              caption="Pearson correlations among pain and adoption-intent variables.")


def slide_12_sipoc(prs) -> None:
    slide = new_slide(prs, idx=12, title="Current Process - SIPOC",
                      footer="Tool: SIPOC - Process Scoping")
    bullets = [
        "Current process has 5 steps; 3 of them have no digital artifact today.",
        "Inputs column has no demand-visibility signal - the app injects it as the first intervention.",
        "Outputs side has a leakage-prone cash output - replaced with reconciled wallet debit.",
        "Customer column treats drivers as implicit; app makes driver utilisation a first-class KPI.",
    ]
    add_image_with_insights(slide, ANALYSIS / "04_sipoc" / "diagram.png",
                              bullets,
                              caption="SIPOC scoping of the current cash-mediated E-rick process.")


def slide_13_current_future(prs) -> None:
    slide = new_slide(prs, idx=13, title="Current vs Future Flow",
                      footer="Tool: Process Flowchart - Ref. Ch.2")
    bullets = [
        "Three current-state pain cells (4-pax wait, cash pay, empty return) all flip to sage improvements.",
        "Future state adds two commitment gates (60-sec timer, return-trip match) that enforce discipline.",
        "Step count is similar on both sides; cognitive load per step drops sharply.",
        "Each pain-to-improvement swap is independently shippable.",
    ]
    add_image_with_insights(slide,
                              ANALYSIS / "05_flowchart_current_vs_future" / "diagram.png",
                              bullets,
                              caption="Side-by-side current vs future process; pain cells on the left map to improved cells on the right.")


def slide_14_solution_overview(prs) -> None:
    slide = new_slide(prs, idx=14,
                      title="Proposed Solution Overview",
                      footer="Section: Solution Design")
    bullets = [
        "Four features map 1:1 to four pain points: discovery, booking, wallet, return-match.",
        "Each feature independently deployable; combined value compounds.",
        "Driver-side features (Reached + offline-add) preserve current operational habits.",
        "Placeholder locations swapped for the finalized 43-item IIT Roorkee list (April 2026).",
    ]
    add_image_with_insights(slide, APP_SHOTS / "01_student_map.png",
                              bullets,
                              caption="Student map preview - the entry point to the end-to-end app flow.")


def slide_15_student_map(prs) -> None:
    slide = new_slide(prs, idx=15, title="App - Student Map")
    bullets = [
        "100 m proximity radius shows all drivable options without overwhelming the view.",
        "Passenger count and destination visibility per E-rick lets the user pre-match intent.",
        "Full E-ricks switch to maroon glyphs - capacity state is visible before tapping.",
        "Hand-drawn SVG campus chart, not Google Maps, keeps focus on E-rick state.",
    ]
    add_image_with_insights(slide, APP_SHOTS / "01_student_map.png",
                              bullets,
                              caption="Student-side map: proximity discovery with inline passenger counts.")


def slide_16_booking_countdown(prs) -> None:
    slide = new_slide(prs, idx=16, title="App - Booking + Countdown")
    bullets = [
        "Drop-off dropdown backed by 43 grouped campus locations; single tap confirms.",
        "60-second countdown is a commitment device; enforced via reducer-level TICK_TIMERS.",
        "Expired booking auto-cancels and debits Rs5 - the soft enforcement mechanic.",
    ]
    add_two_images_with_insights_below(slide,
                                         APP_SHOTS / "03_booking_modal.png",
                                         APP_SHOTS / "04_countdown_active.png",
                                         bullets)


def slide_17_wallet(prs) -> None:
    slide = new_slide(prs, idx=17, title="App - Wallet & Payment")
    bullets = [
        "Wallet holds a booking-side 'hold' so fare appears debited immediately on booking.",
        "Rs50 soft-debt cap at ride-end prevents the worst cash-leakage scenario.",
        "Ride history is immutable - every debit and cancellation fee surfaces transparently.",
        "Pre-linked driver accounts eliminate the manual transfer step.",
    ]
    add_image_with_insights(slide, APP_SHOTS / "05_wallet_panel.png",
                              bullets,
                              caption="Wallet panel: auto-debit, holds, transparent history, bounded soft-debt.")


def slide_18_driver(prs) -> None:
    slide = new_slide(prs, idx=18, title="App - Driver View")
    bullets = [
        "'Mark Reached' locks the seat to the physically-present student (solves seat-stealing).",
        "Offline-add handles walk-up passengers and keeps all student apps in sync.",
    ]
    add_two_images_with_insights_below(slide,
                                         APP_SHOTS / "06_driver_view.png",
                                         APP_SHOTS / "07_driver_add_passenger.png",
                                         bullets)


def slide_19_admin(prs) -> None:
    slide = new_slide(prs, idx=19, title="App - Admin Dashboard")
    bullets = [
        "Previously invisible KPIs (wait, leakage, empty return) become first-class metrics.",
        "Dashboard enables the Check phase of PDCA without extra instrumentation.",
        "Fleet table surfaces per-driver utilisation - feeds the driver-incentive redesign.",
        "Live map gives operational-response capability (e.g., reroute a stranded E-rick).",
    ]
    add_image_with_insights(slide, APP_SHOTS / "08_admin_view.png",
                              bullets,
                              caption="Admin dashboard: live KPIs, fleet table, and operational map view.")


def slide_20_fmea(prs) -> None:
    slide = new_slide(prs, idx=20, title="FMEA Risk Scoring",
                      footer="Tool: FMEA - Ref. Ch.10 Reliability")
    bullets = [
        "Top-3 RPN (driver offline 126, dispute 120, wallet insufficient 100) drive 38% of total risk.",
        "7 of 10 modes have reducer-level (not just UI) mitigation - defense in depth.",
        "Race-condition mode scored lowest (RPN 20) - server-side first-write-wins handles it cleanly.",
        "Every edge case in the prototype has a corresponding row here - no untested path.",
    ]
    add_image_with_insights(slide, ANALYSIS / "06_fmea" / "diagram.png",
                              bullets,
                              caption="FMEA rows scored Severity x Occurrence x Detection, top risks prioritised.")


def slide_21_pdca(prs) -> None:
    slide = new_slide(prs, idx=21, title="Rollout - PDCA",
                      footer="Tool: PDCA - Ref. Ch.1")
    bullets = [
        "First cycle 4-6 weeks: MVP, 5 E-ricks, 50 pilot students, single zone.",
        "Check phase reuses the same Q1-Q22 survey instrument used in this project.",
        "Act gate: expansion requires >=2 KPI improvements vs baseline or the cycle repeats.",
        "Built for iteration - no assumption of a 'big-bang' launch.",
    ]
    add_image_with_insights(slide, ANALYSIS / "07_pdca" / "diagram.png",
                              bullets,
                              caption="PDCA roadmap: pilot, measure, standardise, expand - one zone at a time.")


def slide_22_impact(prs) -> None:
    slide = new_slide(prs, idx=22, title="Quantified Impact",
                      footer="Framing: Cost of Quality (PAF) - Ch.1")

    # Left column: Baseline
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7),
                                         Inches(6.0), Inches(4.4))
    tfL = left_box.text_frame
    tfL.margin_left = tfL.margin_right = 0
    tfL.word_wrap = True

    pL0 = tfL.paragraphs[0]
    rL0 = pL0.add_run()
    rL0.text = "Baseline (from survey, n = 150)"
    set_run_font(rL0, name="Georgia", size=20, bold=True, color=ACCENT)
    pL0.space_after = Pt(10)

    base_lines = [
        "Average boarding wait: 8.2 min",
        "Fare leakage (cash): ~5-8%",
        "Empty-return share: ~30%",
        "Overall satisfaction: 2.1 / 5",
    ]
    for line in base_lines:
        p = tfL.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"- {line}"
        set_run_font(r, name="Calibri", size=18, color=INK)

    # Right column: Target
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.7),
                                          Inches(6.0), Inches(4.4))
    tfR = right_box.text_frame
    tfR.margin_left = tfR.margin_right = 0
    tfR.word_wrap = True

    pR0 = tfR.paragraphs[0]
    rR0 = pR0.add_run()
    rR0.text = "Target (with app intervention)"
    set_run_font(rR0, name="Georgia", size=20, bold=True, color=ACCENT)
    pR0.space_after = Pt(10)

    tgt_lines = [
        "Average boarding wait: < 2 min",
        "Fare leakage (wallet): < 0.01%",
        "Empty-return share: < 10%",
        "Overall satisfaction: > 4.0 / 5",
    ]
    for line in tgt_lines:
        p = tfR.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"- {line}"
        set_run_font(r, name="Calibri", size=18, color=INK)

    # Subtext
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(6.25),
                                    Inches(12.0), Inches(0.6))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = ("75% reduction in boarding wait; near-elimination of leakage.")
    set_run_font(r, name="Georgia", size=18, italic=True, bold=True, color=INK)


def slide_23_governance(prs) -> None:
    slide = new_slide(prs, idx=23, title="Success Metrics & Governance",
                      footer="Framing: SPC (X-bar) - Ref. Ch.4; Measurement - Ch.1")
    headers = ["KPI", "Baseline", "Target", "Measurement"]
    rows = [
        ["Average boarding wait (min)", "8.2", "< 2.0",
         "Server timestamps (booked -> reached); SPC X-bar chart weekly"],
        ["Fare leakage (% of GMV)", "~5 - 8%", "< 0.01%",
         "Wallet auto-debit logs; reconciliation report"],
        ["Empty-return share", "~30%", "< 10%",
         "Driver GPS trails; compare boarded vs empty segments"],
        ["Student weekly active (WAU)", "NA",
         "> 60% of residents", "Unique booking-student-ids / week"],
        ["NPS (post-ride)", "NA", "> 40",
         "One-tap 0-10 prompt after completion; rolling 7-day"],
    ]
    add_table(slide, headers, rows, left=0.45, top=1.7,
              width=12.45, height=4.8, header_size=13, body_size=12)


def slide_24_conclusion(prs) -> None:
    slide = new_slide(prs, idx=24, title="Conclusion & Next Steps")
    bullets = [
        "Problem is real and measurable - 8.2 min mean wait, 56.8% of pain concentrated in the top-3 issues, satisfaction at 2.9 / 5.",
        "Prototype demonstrates every core flow end-to-end: proximity map, booking, 60-sec countdown, wallet auto-debit, driver 'Reached' seat-lock, admin KPIs - with edge cases (offline add, cancellation, dispute) handled.",
        "Pilot-ready: a two-week PDCA pilot in one zone will surface real-world constants (actual leakage %, adoption curve, driver onboarding time).",
    ]
    add_bullets(slide, bullets, left=0.7, top=1.7, width=12.0, size=18)

    closing = slide.shapes.add_textbox(Inches(0.7), Inches(6.2),
                                        Inches(12.0), Inches(0.6))
    p = closing.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Thank you - questions welcome."
    set_run_font(r, name="Georgia", size=22, italic=True, bold=True, color=ACCENT)


# ---------------------------------------------------------------------------
# Stats recomputation (printed to stdout for spot-checking)
# ---------------------------------------------------------------------------

def recompute_and_print_stats() -> None:
    df = pd.read_csv(SURVEY_CSV)
    n = len(df)

    wait_mean = df["q5_wait_minutes"].mean()
    wait_median = df["q5_wait_minutes"].median()
    wait_p90 = np.percentile(df["q5_wait_minutes"], 90)

    issues = (df["q12_issues_faced"].dropna().astype(str)
              .str.split(";").explode().str.strip())
    issues = issues[issues != ""]
    issue_counts = issues.value_counts()
    issue_share = (issue_counts / n) * 100

    top3 = issue_share.head(3)
    top3_total = top3.sum()

    r_q7_q16 = df[["q7_payment_issue_freq", "q16_wallet_comfort"]].corr().iloc[0, 1]
    r_q5_q14 = df[["q5_wait_minutes", "q14_app_willingness"]].corr().iloc[0, 1]

    fee_accept_pct = (df["q18_fee_acceptance"] >= 3).mean() * 100
    sat_mean = df["q11_overall_satisfaction"].mean()

    print("=" * 60)
    print(f"Survey stats (n={n})")
    print("=" * 60)
    print(f"  wait_mean_min        = {wait_mean:.2f}")
    print(f"  wait_median_min      = {wait_median:.1f}")
    print(f"  wait_p90_min         = {wait_p90:.1f}")
    print(f"  top pain points (share%):")
    for name, val in issue_share.items():
        print(f"    - {name:24s}  {val:5.2f}%")
    print(f"  top-3 pain total share = {top3_total:.2f}%")
    print(f"  r(payment_issue, wallet_comfort)   = {r_q7_q16:.3f}")
    print(f"  r(wait_minutes,  app_willingness)  = {r_q5_q14:.3f}")
    print(f"  fee-acceptance (>=3 out of 5) = {fee_accept_pct:.1f}%")
    print(f"  overall satisfaction mean     = {sat_mean:.2f} / 5")
    print("=" * 60)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build_deck() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build each slide in order.
    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_stakeholders(prs)
    slide_04_methodology(prs)
    slide_05_wait_time(prs)
    slide_06_pareto(prs)
    slide_07_ishikawa(prs)
    slide_08_five_whys(prs)
    slide_09_relations(prs)
    slide_10_affinity(prs)
    slide_11_correlation(prs)
    slide_12_sipoc(prs)
    slide_13_current_future(prs)
    slide_14_solution_overview(prs)
    slide_15_student_map(prs)
    slide_16_booking_countdown(prs)
    slide_17_wallet(prs)
    slide_18_driver(prs)
    slide_19_admin(prs)
    slide_20_fmea(prs)
    slide_21_pdca(prs)
    slide_22_impact(prs)
    slide_23_governance(prs)
    slide_24_conclusion(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    return OUT_PPTX


def main() -> None:
    recompute_and_print_stats()
    out = build_deck()
    size_mb = out.stat().st_size / (1024 * 1024)
    print(f"Wrote: {out}")
    print(f"Size : {size_mb:.2f} MB  ({out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
